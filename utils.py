from docx import Document
from pptx import Presentation
import openpyxl
from pdf2image import convert_from_bytes
import pytesseract
import re
import aiohttp
import io
from PIL import Image
#import speech_recognition as sr
#from moviepy.editor import VideoFileClip
import magic

def extract_text_from_docx(docx_bytes):
    doc = Document(io.BytesIO(docx_bytes))
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    return text

def extract_text_from_pptx(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def extract_text_from_excel(attachment_file):
    workbook = openpyxl.load_workbook(filename=attachment_file)
    extracted_text = []

    # Iterate over each sheet in the workbook
    for sheet in workbook.sheetnames:
        # Select the current sheet
        current_sheet = workbook[sheet]
        # Iterate over each row in the sheet
        for row in current_sheet.iter_rows(values_only=True):
            # Iterate over each cell in the row
            for cell in row:
                # Check if the cell contains text
                if cell.value:
                    extracted_text.append(str(cell.value))
    return '\n'.join(extracted_text)

async def extract_text_from_pdf(attachment_file):
    # Convert PDF to images
    images = convert_from_bytes(attachment_file)

    # Initialize an empty string to store the extracted text
    text = ""

    # Iterate over each page image and perform OCR
    for img in images:
        page_text = pytesseract.image_to_string(img)
        text += page_text

    return text

async def extract_text_from_image(attachment_file):
    try:
        # Create a binary stream from the attachment file
        img_stream = io.BytesIO(attachment_file)
        # Open the image from the binary stream
        with Image.open(img_stream) as img:
            # Perform OCR on the image
            text = pytesseract.image_to_string(img)
            return text
    except Exception as e:
        # Handle any exceptions
        print(f"Error: {e}")
        return None

'''
async def extract_text_from_audio(audio_file):
    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_file) as source:
        audio_data = recognizer.record(source)
        text = recognizer.recognize_google(audio_data)
        return text

async def extract_text_from_video(video_file):
    video_clip = VideoFileClip(video_file)
    audio_clip = video_clip.audio
    audio_clip.write_audiofile("temporary_audio.wav")
    text = await extract_text_from_audio("temporary_audio.wav")
    return text
'''

async def download_file(url):
    async with aiohttp.ClientSession() as session:
        async with session.get(url) as response:
            if response.status == 200:
                return await response.read()
            else:
                # Handle error responses here
                raise ValueError(f"Failed to download file: {response.status}")
            
async def get_file_content(attachment):

    try:
        # If attachment is uploaded from Discord, download it
        attachment_url = attachment.url
        attachment_file = await download_file(attachment_url)
    except AttributeError:
        # If attachment is sent as a link, leave it be.
        attachment_file = attachment

    # Determine file type based on magic number
    file_type = magic.from_buffer(attachment_file)
    file_type = file_type.lower()
    print("Detected file type:", file_type)

    if 'text' in file_type:
        # Text file
        text = attachment_file.decode('utf-8')
        return [text, "text"]
    elif 'word' in file_type:
        # Word document
        text = extract_text_from_docx(attachment_file)
        return [text, "document"]
    elif 'powerpoint' in file_type:
        # PowerPoint document
        text = extract_text_from_pptx(attachment_file)
        return [text, "PowerPoint"]
    elif 'excel' in file_type:
        # Excel document
        text = extract_text_from_excel(attachment_file)
        return [text, "data"]
    elif 'pdf' in file_type:
        # PDF file
        text = await extract_text_from_pdf(attachment_file)
        return [text, "text"]
    elif 'image' in file_type:
        # Image file
        text = await extract_text_from_image(attachment_file)
        return [text, "text"]
    else:
        return False
    '''
    elif 'audio' in file_type or any(ext in file_type for ext in ['wav', 'mp3', 'aac']):
        text = await extract_text_from_audio(attachment_directory)
        return [text, "audio"]
    elif 'video' in file_type or any(ext in file_type.encode() for ext in [b'.mp4', b'.avi', b'.mov', b'.mkv']):
        # Check if the filename ends with any of the specified video extensions
        text = await extract_text_from_video(attachment_directory)
    '''
    

def extract_video_id(url):
    match = re.search(r'(?:https?://)?(?:www\.)?(?:youtube\.com/watch\?v=|youtu\.be/)([a-zA-Z0-9_-]{11})', url)
    if match:
        return match.group(1)
    else:
        return None
    
def get_length(word_count):
    if word_count < 150:
        return "very short"
    elif word_count < 200:
        return "short"
    elif word_count < 300:
        return "long"
    else:
        return "very long"