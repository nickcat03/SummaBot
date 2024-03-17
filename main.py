# This example requires the 'message_content' privileged intents

import requests
import json
import discord
import os
from os import getenv
from discord import app_commands
from discord.ext import commands
from discord.ext.commands import cooldown, BucketType
from discord.app_commands import Choice
import re
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import io


def extract_text_from_docx(docx_bytes):
    doc = Document(io.BytesIO(docx_bytes))
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    return text

def extract_text_from_pptx(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text


#KEYS = json.load(open("keys.json"))

# Initialize Discord bot intents with default settings, including basic events like join/leave.
intents = discord.Intents.default()

# Enable the bot to read message content, necessary for responding to user messages.
intents.message_content = True

# Create a bot instance with a command prefix '!', ready to handle commands and intents.
client = commands.Bot(command_prefix="!", intents=intents)


# # This is a slash command that calls the GPT-3 API to generate a response and sends it back to the user on discord
# @client.tree.command(name='query_gpt', description='Ask GPT a question')
# @app_commands.describe(query = "Ask GPT a question")
# @app_commands.checks.cooldown(1, 60, key=lambda i: (i.guild_id, i.user.id)) # This prevents people from spamming the command. They can only use it once every 60 seconds.
# async def query_gpt(interaction: discord.Interaction, query: str):
    
#     completion = openai_client.chat.completions.create(
#         model="gpt-3.5-turbo",
#         messages = [
#             {"role": "system", "content": "You are a helpful assistant"},
#             {"role": "user", "content": query}
#         ],
#     )
    
#     response = completion.choices[0].message.content
#     await interaction.response.send_message(response)
    

# Helper function to encode images to base64.
# Used to send images to the GPT-4 Vision API.
# def encode_image(image_path):
#     with open(image_path, "rb") as image_file:
#         return base64.b64encode(image_file.read()).decode('utf-8')

@client.event
async def on_message(message):
    # If bot sends a command, don't run it as itself

    # Set default number of words
    num_of_words = 400

    if message.author == client.user:
        return
  
    # Generate summary from text
    if message.content.lower().startswith('summarize'):
        split_message = message.content.split()
        input_text = ''

        # Check for attachments
        if message.attachments:
            # Assuming only one attachment is allowed
            attachment = message.attachments[0]
            file_extension = attachment.filename.split('.')[-1].lower()

            if file_extension == 'txt':
                # Text file
                attachment_content = await attachment.read()
                input_text = attachment_content.decode('utf-8')
            elif file_extension in ['doc', 'docx']:
                # Word document
                attachment_content = await attachment.read()
                input_text = extract_text_from_docx(attachment_content)
            elif file_extension == 'pptx':
                # PowerPoint document
                attachment_content = await attachment.read()
                input_text = extract_text_from_pptx(attachment_content)
            else:
                await message.channel.send("Invalid attachment format. Please attach a .txt, .doc, .docx, or .pptx file.")
                return

        else:
            input_text = ' '.join(split_message[2:])  # Extract the text or URL from the message

        # Check if the user has provided the number of words
        try:
            if len(split_message) > 1:
                num_of_words = int(split_message[1])
                if num_of_words < 100 or num_of_words > 1000:
                    await message.channel.send("Number of words should be between 100 and 1000.")
                    return
        except ValueError:
            await message.channel.send("Invalid command format. Please use 'summarize [num_of_words] (text or URL)' to generate summaries.")
            return

        # Check if the input is a URL
        if re.match(r'https?://\S+', input_text):
            # If it's a URL, fetch content from the URL
            try:
                response = requests.get(input_text)
                response.raise_for_status()  # Raise an exception for HTTP errors
                html_content = response.text
                soup = BeautifulSoup(html_content, 'html.parser')
                text_content = soup.get_text()
            except requests.exceptions.RequestException as e:
                await message.channel.send(f"Error: {e}")
                return
        else:
            # If it's not a URL, use the input text directly
            text_content = input_text

        print(text_content)
        payload = {
            "model": "google/gemma-7b-it:free",
            "messages": [
                {
                    "role": "user",
                    "content": f"You are a journalist. Summarize the following, make sure the summary is concise. Only write {num_of_words} words, and do not exceed 2000 characters. Do not give me other information outside of those paragraphs: {text_content}."
                }
            ]
        }

        try:
            response = requests.post(
                url="https://openrouter.ai/api/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {os.environ["OPEN_BORDER_KEY"]}",
                },
                data=json.dumps(payload)
            )

            summary = response.json()['choices'][0]['message']['content']
            print("SUMMARY:", summary)
            
            await message.channel.send(summary)

        except Exception as e:
            await message.channel.send(f"An error occurred: {str(e)}")
        

    #generate questions from text
    if message.content.lower().startswith('questions'):
        # Split the message content by spaces
        split_message = message.content.split()
        
        # Extract the number of questions and text content from the message
        try:
            num_questions = int(split_message[1])  # Extract the number from the message
            input_text = ' '.join(split_message[2:])  # Extract the text content after the number
        except (ValueError, IndexError):
            await message.channel.send("Invalid command format. Please use 'questions [number] (text)' to generate questions.")
            return
        
        # Check for attachments
        if message.attachments:
            # Assuming only one attachment is allowed
            attachment = message.attachments[0]
            file_extension = attachment.filename.split('.')[-1].lower()

            if file_extension == 'txt':
                # Text file
                attachment_content = await attachment.read()
                input_text = attachment_content.decode('utf-8')
            elif file_extension in ['doc', 'docx']:
                # Word document
                attachment_content = await attachment.read()
                input_text = extract_text_from_docx(attachment_content)
            elif file_extension == 'pptx':
                # PowerPoint document
                attachment_content = await attachment.read()
                input_text = extract_text_from_pptx(attachment_content)
            else:
                await message.channel.send("Invalid attachment format. Please attach a .txt, .doc, .docx, or .pptx file.")
                return

        else:
            input_text = ' '.join(split_message[2:])  # Extract the text or URL from the message
        
        if re.match(r'https?://\S+', input_text):
            # If it's a URL, fetch content from the URL
            try:
                response = requests.get(input_text)
                response.raise_for_status()  # Raise an exception for HTTP errors
                html_content = response.text
            except requests.exceptions.RequestException as e:
                await message.channel.send(f"Error: {e}")
                return

            # Parse HTML content to extract text
            soup = BeautifulSoup(html_content, 'html.parser')
            text_content = soup.get_text()
        else:
            # If it's not a URL, use the input text directly
            text_content = input_text

        payload = {
            "model": "google/gemma-7b-it:free",
            "messages": [
                {
                    "role": "user",
                    "content": f"You are a helpful assistant. Generate {num_questions} questions based on the following. Do not use any other reference, only utilize the text given here: {text_content}. When generaating your response, do not write anything else. Only send the questions. When generating questions, space them out with one single line break, do not use multiple. Generate {num_questions} questions."
                }
            ]
        }

        try:
            response = requests.post(
                url="https://openrouter.ai/api/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {os.environ["OPEN_BORDER_KEY"]}",
                },
                data=json.dumps(payload)
            )

            questions = response.json()['choices'][0]['message']['content']
            print("SUMMARY:", questions)
            
            await message.channel.send(questions)

        except Exception as e:
            print("Error:", e)
            await message.channel.send("An error has occured.")
            return

                
    # Test command to respond with "hello" when user types "hi"
    if message.content.lower() == 'hi':
        await message.channel.send('helloooooooo :)')




# Decorator to register an error handler for application command errors in the bot
@client.tree.error
async def on_app_command_error(interaction: discord.Interaction, error: app_commands.AppCommandError):
    # Check if the error is a command cooldown error
    if isinstance(error, app_commands.CommandOnCooldown):
        # Send an ephemeral message informing the user of the cooldown with the remaining time
        await interaction.response.send_message(f"Please wait for {round(error.retry_after, 0)} seconds before running the command again!", ephemeral=True)
    else:
        # [Optional: Modify this part to handle other types of errors appropriately]
        # Currently, sends the same cooldown message for all types of errors, which may not be intended
        await interaction.response.send_message(f"Please wait for {round(error.retry_after, 0)} seconds before running the command again!", ephemeral=True)



# This is the code that runs the bot.
# Make sure to replace the DISCORD_TOKEN below with your bot's token.
# Get your bot token from the discord developer portal: https://discordapp.com/developers/applications/
# Use these articles to help you get started:
#  https://www.writebots.com/discord-bot-token/
#  https://realpython.com/how-to-make-a-discord-bot-python/

client.run(os.environ["DISCORD_TOKEN"])
